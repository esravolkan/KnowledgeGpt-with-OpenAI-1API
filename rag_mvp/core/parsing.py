from io import BytesIO
from typing import List, Any, Optional
import re
from dataclasses import dataclass

import docx2txt
from langchain.docstore.document import Document
import fitz
from hashlib import md5
from xlsx2html import xlsx2html
from pptx import Presentation
from openpyxl import load_workbook

from abc import abstractmethod, ABC
from copy import deepcopy


@dataclass(init=False)
class File(ABC):
    """Represents an uploaded file comprised of Documents"""

    def __init__(
        self,
        name: str,
        id: str,
        metadata: Optional[dict[str, Any]] = None,
        docs: Optional[List[Document]] = None,
    ):
        self.name = name
        self.id = id
        self.metadata = metadata or {}
        self.docs = docs or []

    @classmethod
    @abstractmethod
    def from_bytes(cls, file: BytesIO) -> "File":
        """Creates a File from a BytesIO object"""

    def __repr__(self) -> str:
        return (
            f"File(name={self.name}, id={self.id},"
            " metadata={self.metadata}, docs={self.docs})"
        )

    def __str__(self) -> str:
        return f"File(name={self.name}, id={self.id}, metadata={self.metadata})"

    def copy(self) -> "File":
        """Create a deep copy of this File"""
        return self.__class__(
            name=self.name,
            id=self.id,
            metadata=deepcopy(self.metadata),
            docs=deepcopy(self.docs),
        )


def strip_consecutive_newlines(text: str) -> str:
    """Strips consecutive newlines from a string
    possibly with whitespace in between
    """
    return re.sub(r"\s*\n\s*", "\n", text)


class DocxFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "DocxFile":
        text = docx2txt.process(file)
        text = strip_consecutive_newlines(text)
        doc = Document(page_content=text.strip())
        doc.metadata["source"] = "p-1"
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=[doc])


class PdfFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "PdfFile":
        pdf = fitz.open(stream=file.read(), filetype="pdf")  # type: ignore
        docs = []
        for i, page in enumerate(pdf):
            text = page.get_text(sort=True)
            text = strip_consecutive_newlines(text)
            doc = Document(page_content=text.strip())
            doc.metadata["page"] = i + 1
            doc.metadata["source"] = f"p-{i+1}"
            docs.append(doc)
        # file.read() mutates the file object, which can affect caching
        # so we need to reset the file pointer to the beginning
        file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


class TxtFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "TxtFile":
        text = file.read().decode("utf-8", errors="replace")
        text = strip_consecutive_newlines(text)
        file.seek(0)
        doc = Document(page_content=text.strip())
        doc.metadata["source"] = "p-1"
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=[doc])


class XlsmFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "XlsmFile":
        workbook = load_workbook(file)
        docs = []

        for i, sheet in enumerate(workbook.sheetnames):
            html = xlsx2html(file, sheet=sheet) #f'{sheet}.html',
            html.seek(0)
            doc = Document(page_content=html.read().strip())
            html.seek(0)
            doc.metadata["page"] = i + 1
            doc.metadata["source"] = f"p-{sheet}"
            docs.append(doc)
            #html.seek(0)
        file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


class PptxFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "PptxFile":
        presentation = Presentation(file)
        docs = []
        for slide_number, slide in enumerate(presentation.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            text = strip_consecutive_newlines(text)
            doc = Document(page_content=text.strip())
            doc.metadata["page"] = slide_number + 1
            doc.metadata["source"] = f"p-{slide_number + 1}"
            docs.append(doc)
            file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


def read_file(file: BytesIO) -> File:
    """Reads an uploaded file and returns a File object"""
    if file.name.lower().endswith(".docx"):
        return DocxFile.from_bytes(file)
    elif file.name.lower().endswith(".pdf"):
        return PdfFile.from_bytes(file)
    elif file.name.lower().endswith(".txt"):
        return TxtFile.from_bytes(file)
    elif file.name.lower().endswith(".xlsx"):
        return XlsmFile.from_bytes(file)
    elif file.name.lower().endswith(".pptx"):
        return PptxFile.from_bytes(file)
    else:
        raise NotImplementedError(f"File type {file.name.split('.')[-1]} not supported")
